"""
文档解析服务
支持 PDF / Word / Excel / PPT / TXT / 图片
纯文本提取 + 分块
"""
import re
import tempfile
from typing import Optional
from pathlib import Path

from backend.utils.logger import logger
from backend.services.advanced_document_parser import AdvancedDocumentParser, HAS_FITZ

# 尝试导入文档解析库（非强制，缺失时给出清晰错误）
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class ParseResult:
    """文档解析结果"""
    def __init__(self, text: str, chunks: list[dict], page_count: int = 0):
        self.text = text          # 完整文本
        self.chunks = chunks      # [{content, metadata}]
        self.page_count = page_count


class DocumentParser:
    """多格式文档解析器"""

    CHUNK_SIZE = 500      # 每块字符数
    CHUNK_OVERLAP = 50    # 块重叠

    def __init__(self):
        self.advanced = AdvancedDocumentParser(self._chunk_text)

    def parse(self, file_path: str | Path, file_type: str = "", chunk_template: str = "general") -> Optional[ParseResult]:
        """
        解析文档，返回结构化结果
        file_type: pdf|docx|xlsx|pptx|txt|image
        """
        path = Path(file_path)
        if not path.exists():
            logger.error(f"File not found: {file_path}")
            return None

        # 推断文件类型
        if not file_type:
            ext = path.suffix.lower()
            type_map = {
                ".pdf": "pdf", ".docx": "docx", ".doc": "docx",
                ".xlsx": "xlsx", ".xls": "xlsx",
                ".pptx": "pptx", ".ppt": "pptx",
                ".txt": "txt", ".md": "txt", ".csv": "txt",
                ".jpg": "image", ".jpeg": "image", ".png": "image",
                ".gif": "image", ".bmp": "image", ".tiff": "image",
            }
            file_type = type_map.get(ext, "txt")

        try:
            if file_type == "pdf":
                if HAS_FITZ:
                    return self.advanced.parse_pdf(path, chunk_template, ParseResult)
                return self._parse_pdf(path, chunk_template)
            elif file_type == "docx":
                advanced = self.advanced.parse_docx(path, chunk_template, ParseResult)
                return advanced or self._parse_docx(path, chunk_template)
            elif file_type == "xlsx":
                return self._parse_xlsx(path, chunk_template)
            elif file_type == "pptx":
                return self._parse_pptx(path, chunk_template)
            elif file_type == "image":
                return self._parse_image(path, chunk_template)
            else:  # txt / 其他
                return self._parse_txt(path, chunk_template)
        except Exception as e:
            logger.error(f"Parse {file_type} error: {e}", exc_info=True)
            return None

    def _chunk_text(
        self,
        text: str,
        metadata: Optional[dict] = None,
        template: str = "general",
    ) -> list[dict]:
        """Split text with an explainable template and exact character overlap."""
        if not text or not text.strip():
            return []

        configs = {
            "general": (600, 80),
            "sentence": (350, 50),
            "table": (1200, 0),
            "qa": (800, 60),
        }
        size, overlap = configs.get(template, configs["general"])
        clean_text = text.replace("\r\n", "\n").strip()

        if template == "table":
            units = [line.strip() for line in clean_text.splitlines() if line.strip()]
        elif template == "sentence":
            units = [item.strip() for item in re.split(r"(?<=[。！？!?；;])\s*", clean_text) if item.strip()]
        elif template == "qa":
            units = [item.strip() for item in re.split(r"\n\s*\n|(?=^(?:问|Q|问题)[:：])", clean_text, flags=re.M) if item.strip()]
        else:
            units = [item.strip() for item in re.split(r"\n\s*\n", clean_text) if item.strip()]

        normalized_units: list[str] = []
        for unit in units:
            if len(unit) <= size:
                normalized_units.append(unit)
                continue
            sentences = [item.strip() for item in re.split(r"(?<=[。！？!?；;])\s*", unit) if item.strip()]
            for sentence in sentences:
                if len(sentence) <= size:
                    normalized_units.append(sentence)
                else:
                    normalized_units.extend(sentence[start:start + size] for start in range(0, len(sentence), size))

        chunks: list[dict] = []
        current = ""
        cursor = 0
        for unit in normalized_units:
            candidate = f"{current}\n{unit}".strip() if current else unit
            if current and len(candidate) > size:
                chunk_meta = dict(metadata or {})
                chunk_meta.update({
                    "chunk_index": len(chunks), "template": template,
                    "char_start": cursor, "char_end": cursor + len(current),
                })
                chunks.append({"content": current, "metadata": chunk_meta})
                cursor += max(1, len(current) - overlap)
                prefix = current[-overlap:] if overlap else ""
                current = f"{prefix}\n{unit}".strip() if prefix else unit
            else:
                current = candidate

        if current:
            chunk_meta = dict(metadata or {})
            chunk_meta.update({
                "chunk_index": len(chunks), "template": template,
                "char_start": cursor, "char_end": cursor + len(current),
            })
            chunks.append({"content": current, "metadata": chunk_meta})
        return chunks
    # ---------- PDF ----------

    def _parse_pdf(self, path: Path, template: str = "general") -> Optional[ParseResult]:
        if not HAS_PDF:
            logger.warning("PyPDF2 not installed, cannot parse PDF")
            return None

        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            pages = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                pages.append(text)

            full_text = "\n\n".join(pages)
            chunks = []
            for page_number, page_text in enumerate(pages, 1):
                chunks.extend(self._chunk_text(page_text, {"page": page_number}, template))
            return ParseResult(text=full_text, chunks=chunks, page_count=len(pages))

    # ---------- Word ----------

    def _parse_docx(self, path: Path, template: str = "general") -> Optional[ParseResult]:
        if not HAS_DOCX:
            logger.warning("python-docx not installed, cannot parse DOCX")
            return None

        doc = docx.Document(str(path))
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())

        full_text = "\n".join(paragraphs)
        return ParseResult(
            text=full_text,
            chunks=self._chunk_text(full_text, template=template),
            page_count=len(paragraphs) // 20 + 1,
        )

    # ---------- Excel ----------

    def _parse_xlsx(self, path: Path, template: str = "general") -> Optional[ParseResult]:
        if not HAS_XLSX:
            logger.warning("openpyxl not installed, cannot parse XLSX")
            return None

        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        all_texts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                text_row = [str(cell) for cell in row if cell is not None]
                if text_row:
                    rows.append(" | ".join(text_row))
            if rows:
                all_texts.append(f"【工作表: {sheet_name}】\n" + "\n".join(rows))

        full_text = "\n\n".join(all_texts)
        chunks = []
        for sheet_text in all_texts:
            sheet_name = sheet_text.split("\n", 1)[0].replace("【工作表: ", "").rstrip("】")
            chunks.extend(self._chunk_text(sheet_text, {"sheet": sheet_name}, template))
        return ParseResult(
            text=full_text,
            chunks=chunks,
            page_count=len(wb.sheetnames),
        )

    # ---------- PPT ----------

    def _parse_pptx(self, path: Path, template: str = "general") -> Optional[ParseResult]:
        if not HAS_PPTX:
            logger.warning("python-pptx not installed, cannot parse PPTX")
            return None

        prs = Presentation(str(path))
        slides = []
        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text.strip())
            if slide_texts:
                slides.append("\n".join(slide_texts))

        full_text = "\n\n---\n\n".join(slides)
        chunks = []
        for slide_number, slide_text in enumerate(slides, 1):
            chunks.extend(self._chunk_text(slide_text, {"page": slide_number, "slide": slide_number}, template))
        return ParseResult(
            text=full_text,
            chunks=chunks,
            page_count=len(slides),
        )

    # ---------- TXT ----------

    def _parse_txt(self, path: Path, template: str = "general") -> Optional[ParseResult]:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            full_text = f.read()

        return ParseResult(
            text=full_text,
            chunks=self._chunk_text(full_text, template=template),
            page_count=1,
        )

    # ---------- 图片（PaddleOCR 识别）----------

    def _parse_image(self, path: Path, template: str = "general") -> Optional[ParseResult]:
        """Defer OCR/vision analysis to the async indexing stage."""
        caption = f"待分析图片: {path.name}"
        return ParseResult(
            text=caption,
            chunks=[{
                "content": caption,
                "metadata": {
                    "image_path": str(path), "file_name": path.name, "page": 1,
                    "content_type": "image", "vision_pending": True,
                },
            }],
            page_count=1,
        )

# 全局单例
document_parser = DocumentParser()
